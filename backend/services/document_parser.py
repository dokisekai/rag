import io
import re
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Dict, Any, Optional

try:
    import chardet
except ImportError:
    chardet = None


@dataclass
class DocumentNode:
    content: str
    node_type: str  # "paragraph", "table", "heading", "list", "slide", "page"
    level: int = 0  # heading level / slide number / page number
    metadata: Dict[str, Any] = field(default_factory=dict)


class BaseParser(ABC):
    supported_extensions: List[str] = []

    @abstractmethod
    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentNode]:
        """返回结构化文档节点列表"""
        pass

    def _clean_text(self, text: str) -> str:
        """通用文本清洗"""
        if not text:
            return ""
        # 合并多个空白行
        text = re.sub(r'\n{3,}', '\n\n', text)
        # 去除行首行尾空白
        lines = [line.strip() for line in text.split('\n')]
        return '\n'.join(lines)


class PDFParser(BaseParser):
    supported_extensions = [".pdf"]

    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentNode]:
        try:
            import fitz
        except ImportError:
            raise RuntimeError("PyMuPDF (fitz) is required for PDF parsing. Install: pip install pymupdf")

        doc = fitz.open(stream=file_bytes, filetype="pdf")
        nodes: List[DocumentNode] = []

        for page_num, page in enumerate(doc, start=1):
            text = page.get_text()
            if not text.strip():
                continue

            # 简单启发式：按空行分段
            paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
            for para in paragraphs:
                # 判断是否为标题：短文本、不以标点结尾
                is_heading = len(para) < 100 and not re.search(r'[。，；：！？\.\,;:!?]$', para)
                node_type = "heading" if is_heading else "paragraph"
                level = 1 if is_heading else 0
                nodes.append(DocumentNode(
                    content=para,
                    node_type=node_type,
                    level=level,
                    metadata={"page": page_num, "source": filename}
                ))

        doc.close()
        return nodes


class DOCXParser(BaseParser):
    supported_extensions = [".docx"]

    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentNode]:
        try:
            import docx
        except ImportError:
            raise RuntimeError("python-docx is required for DOCX parsing. Install: pip install python-docx")

        document = docx.Document(io.BytesIO(file_bytes))
        nodes: List[DocumentNode] = []

        for para in document.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name.lower() if para.style else ""
            is_heading = style_name.startswith("heading")
            level = int(style_name.replace("heading", "").strip()) if is_heading and style_name.replace("heading", "").strip().isdigit() else 0

            nodes.append(DocumentNode(
                content=text,
                node_type="heading" if is_heading else "paragraph",
                level=level,
                metadata={"style": para.style.name if para.style else "", "source": filename}
            ))

        # 提取表格
        for table_idx, table in enumerate(document.tables, start=1):
            rows_text = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                rows_text.append(" | ".join(row_text))
            table_text = "\n".join(rows_text)
            if table_text.strip():
                nodes.append(DocumentNode(
                    content=table_text,
                    node_type="table",
                    level=0,
                    metadata={"table_index": table_idx, "source": filename}
                ))

        return nodes


class PPTXParser(BaseParser):
    supported_extensions = [".pptx"]

    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentNode]:
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError("python-pptx is required for PPTX parsing. Install: pip install python-pptx")

        prs = Presentation(io.BytesIO(file_bytes))
        nodes: List[DocumentNode] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            if slide_texts:
                content = "\n".join(slide_texts)
                nodes.append(DocumentNode(
                    content=content,
                    node_type="slide",
                    level=slide_num,
                    metadata={"slide": slide_num, "source": filename}
                ))

            # 提取备注
            if slide.has_notes_slide and slide.notes_slide:
                notes_text_frame = slide.notes_slide.notes_text_frame
                if notes_text_frame and notes_text_frame.text.strip():
                    nodes.append(DocumentNode(
                        content=notes_text_frame.text.strip(),
                        node_type="paragraph",
                        level=0,
                        metadata={"slide": slide_num, "type": "notes", "source": filename}
                    ))

        return nodes


class MDParser(BaseParser):
    supported_extensions = [".md", ".markdown"]

    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentNode]:
        text = file_bytes.decode("utf-8", errors="ignore")
        text = self._clean_text(text)
        nodes: List[DocumentNode] = []

        for line in text.split('\n'):
            line = line.strip()
            if not line:
                continue

            heading_match = re.match(r'^(#{1,6})\s+(.+)$', line)
            if heading_match:
                level = len(heading_match.group(1))
                content = heading_match.group(2).strip()
                nodes.append(DocumentNode(
                    content=content,
                    node_type="heading",
                    level=level,
                    metadata={"source": filename}
                ))
            else:
                # 合并连续的非标题行为段落
                if nodes and nodes[-1].node_type == "paragraph":
                    nodes[-1].content += "\n" + line
                else:
                    nodes.append(DocumentNode(
                        content=line,
                        node_type="paragraph",
                        level=0,
                        metadata={"source": filename}
                    ))

        # 清洗段落内容
        for node in nodes:
            if node.node_type == "paragraph":
                node.content = re.sub(r'\n{2,}', '\n', node.content).strip()

        return nodes


class TXTParser(BaseParser):
    supported_extensions = [".txt"]

    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentNode]:
        # 编码探测
        encoding = "utf-8"
        if chardet:
            detected = chardet.detect(file_bytes)
            if detected and detected.get("encoding"):
                encoding = detected["encoding"]

        text = file_bytes.decode(encoding, errors="ignore")
        text = self._clean_text(text)
        nodes: List[DocumentNode] = []

        # 按段落切分
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        for para in paragraphs:
            nodes.append(DocumentNode(
                content=para,
                node_type="paragraph",
                level=0,
                metadata={"source": filename}
            ))

        return nodes


class DocumentParserRegistry:
    _parsers: Dict[str, BaseParser] = {}
    _initialized: bool = False

    @classmethod
    def _init_defaults(cls):
        if cls._initialized:
            return
        cls.register(PDFParser())
        cls.register(DOCXParser())
        cls.register(PPTXParser())
        cls.register(MDParser())
        cls.register(TXTParser())
        cls._initialized = True

    @classmethod
    def register(cls, parser: BaseParser):
        for ext in parser.supported_extensions:
            cls._parsers[ext.lower()] = parser

    @classmethod
    def parse(cls, file_bytes: bytes, filename: str) -> List[DocumentNode]:
        cls._init_defaults()
        ext = Path(filename).suffix.lower()
        parser = cls._parsers.get(ext)
        if not parser:
            raise ValueError(f"Unsupported file format: {ext}. Supported: {list(cls._parsers.keys())}")
        return parser.parse(file_bytes, filename)

    @classmethod
    def supported_formats(cls) -> List[str]:
        cls._init_defaults()
        return list(cls._parsers.keys())


# ---------------- Chunking Engine ----------------

@dataclass
class Chunk:
    chunk_id: str
    doc_id: str
    kb_id: str
    content: str
    index_in_doc: int
    heading_path: str = ""
    node_type: str = "paragraph"
    metadata: Dict[str, Any] = field(default_factory=dict)


def get_sentence_overlap(text: str, target_overlap: int = 80) -> str:
    """获取完整的句级/行级重叠切片，防止在单词、加粗标签(**)或半句话中途硬截断"""
    if not text or len(text) <= target_overlap:
        return ""

    tail = text[-min(len(text), target_overlap * 2):]
    delimiters = ['\n\n', '\n', '。', '！', '？', '；', ';', '. ']
    for d in delimiters:
        pos = tail.rfind(d)
        if pos != -1 and pos < len(tail) - 5:
            candidate = tail[pos + len(d):].strip()
            if candidate and len(candidate) <= target_overlap * 1.8:
                return candidate

    lines = [l.strip() for l in tail.split('\n') if l.strip()]
    if len(lines) > 1:
        return lines[-1]

    return ""


class ChunkingEngine:
    def __init__(self, min_len: int = 300, max_len: int = 800, overlap: int = 80):
        self.min_len = min_len
        self.max_len = max_len
        self.overlap = overlap

    def chunk(self, nodes: List[DocumentNode], doc_id: str, kb_id: str) -> List[Chunk]:
        """基于文档节点的结构感知分块"""
        chunks: List[Chunk] = []
        current_heading_path: List[str] = []
        buffer = ""
        buffer_types: List[str] = []
        chunk_index = 0

        def flush_buffer():
            nonlocal buffer, buffer_types, chunk_index
            if not buffer.strip():
                return
            heading_path = " > ".join(current_heading_path) if current_heading_path else ""
            chunks.append(Chunk(
                chunk_id=f"{doc_id}_{chunk_index}",
                doc_id=doc_id,
                kb_id=kb_id,
                content=buffer.strip(),
                index_in_doc=chunk_index,
                heading_path=heading_path,
                node_type=buffer_types[0] if buffer_types else "paragraph",
                metadata={}
            ))
            chunk_index += 1
            buffer = ""
            buffer_types = []

        for node in nodes:
            if node.node_type == "heading":
                # 标题触发缓冲区刷新，并将标题作为后续段落的前缀，不再生成孤立无正文的标题切片
                flush_buffer()
                level = node.level if node.level > 0 else 1
                while len(current_heading_path) >= level:
                    current_heading_path.pop()
                current_heading_path.append(node.content)
                buffer = f"### {node.content}"
                buffer_types = ["heading"]
                continue

            if node.node_type == "table":
                # 表格独立成块
                flush_buffer()
                chunks.append(Chunk(
                    chunk_id=f"{doc_id}_{chunk_index}",
                    doc_id=doc_id,
                    kb_id=kb_id,
                    content=node.content,
                    index_in_doc=chunk_index,
                    heading_path=" > ".join(current_heading_path),
                    node_type="table",
                    metadata=dict(node.metadata)
                ))
                chunk_index += 1
                continue

            if node.node_type in ("paragraph", "slide", "page"):
                # 尝试追加到缓冲区
                proposed = buffer + "\n\n" + node.content if buffer else node.content
                if len(proposed) > self.max_len and buffer:
                    # 保留完整句级/行级重叠，绝不在加粗标记(**)或半句话中途硬截断
                    overlap_text = get_sentence_overlap(buffer, self.overlap)
                    flush_buffer()
                    buffer = overlap_text + "\n\n" + node.content if overlap_text else node.content
                else:
                    buffer = proposed
                if node.node_type not in buffer_types:
                    buffer_types.append(node.node_type)

        flush_buffer()

        # 后处理：合并过短的chunk（除了heading和table）
        return self._merge_short_chunks(chunks)

    def _merge_short_chunks(self, chunks: List[Chunk]) -> List[Chunk]:
        if not chunks:
            return chunks

        merged: List[Chunk] = []
        i = 0
        while i < len(chunks):
            chunk = chunks[i]
            if chunk.node_type in ("heading", "table") or len(chunk.content) >= self.min_len:
                merged.append(chunk)
                i += 1
                continue

            # 尝试向后合并
            combined = chunk.content
            j = i + 1
            while j < len(chunks) and len(combined) < self.min_len:
                if chunks[j].node_type == "table":
                    break
                combined += "\n\n" + chunks[j].content
                j += 1

            if j > i + 1:
                merged.append(Chunk(
                    chunk_id=chunk.chunk_id,
                    doc_id=chunk.doc_id,
                    kb_id=chunk.kb_id,
                    content=combined,
                    index_in_doc=chunk.index_in_doc,
                    heading_path=chunk.heading_path,
                    node_type="paragraph",
                    metadata=chunk.metadata
                ))
                i = j
            else:
                merged.append(chunk)
                i += 1

        return merged

    def chunk_by_sliding_window(self, text: str, doc_id: str, kb_id: str) -> List[Chunk]:
        """纯滑动窗口分块（fallback）"""
        chunks = []
        start = 0
        text_len = len(text)
        index = 0
        while start < text_len:
            end = min(start + self.max_len, text_len)
            chunk_text = text[start:end]
            if len(chunk_text) < self.min_len and start != 0:
                chunk_text = text[start:]
                end = text_len
            chunks.append(Chunk(
                chunk_id=f"{doc_id}_{index}",
                doc_id=doc_id,
                kb_id=kb_id,
                content=chunk_text,
                index_in_doc=index,
                heading_path="",
                node_type="paragraph",
                metadata={}
            ))
            index += 1
            if end >= text_len:
                break
            start = end - self.overlap
        return chunks
